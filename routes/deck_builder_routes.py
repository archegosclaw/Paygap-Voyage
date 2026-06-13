# routes/deck_builder_routes.py
# Deck Builder + Brand Kit API
#
# Endpoints:
#   POST /api/deck/extract-pdf     — PDF → extracted text
#   POST /api/deck/generate        — message + deck JSON → slide operations (LLM)
#   POST /api/deck/export-pptx     — deck JSON → .pptx download
#   GET  /api/brand/kit            — brand.json + asset list
#   POST /api/brand/upload         — upload logo or illustration
#   DELETE /api/brand/asset/{brand}/{category}/{name} — remove asset
#   PUT  /api/brand/palette        — update palette in brand.json

import io
import json
import logging
import os
import re
import tempfile
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, HTTPException, Request, UploadFile, File, Form
from fastapi.responses import JSONResponse, StreamingResponse, Response
from pydantic import BaseModel

from src.auth_helpers import require_user
from src.settings import get_setting

logger = logging.getLogger(__name__)

# ── Paths ──────────────────────────────────────────────────────────────────────
_BASE = Path(__file__).parent.parent
_BRAND_DIR = _BASE / "data" / "brand"
_SKILL_DIR = _BASE / "data" / "skills" / "general" / "mpg-slides"
_ALLOWED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/svg+xml", "image/webp"}

# ── Dynamic brand assets ───────────────────────────────────────────────────────
def _brand_assets_context() -> str:
    """Return a string listing available illustrations and logos for the system prompt."""
    lines = []
    for category, label in (("illustrations", "Illustrations"), ("logos", "Logos")):
        cat_dir = _BRAND_DIR / "mpg" / category
        if not cat_dir.exists():
            continue
        # De-duplicate stems (may have both .svg and .png)
        seen = set()
        names = []
        for f in sorted(cat_dir.iterdir()):
            if f.is_file() and f.suffix.lower() in {".png", ".svg", ".jpg", ".webp"}:
                if f.stem not in seen:
                    seen.add(f.stem)
                    names.append(f.stem)
        if names:
            lines.append(f"{label} ({len(names)}): " + ", ".join(names))
    return "\n".join(lines)


# ── System prompt ─────────────────────────────────────────────────────────────
def _load_system_prompt() -> str:
    brand_assets = _brand_assets_context()

    return f"""You are a slide deck builder for MalaysianPAYGAP (MPG). Generate structured JSON operations to create or edit decks.

## Brand System

**Colours**
- MPG Purple `#B8B4FF` — accent, headings, left card borders, bullet arrows
- MPG Purple Dark `#7B75E8` — strong emphasis
- MPG Purple Light `#E8E6FF` — card backgrounds, right title panel
- Background `#F2F2F2` — all slide backgrounds
- Surface `#FFFFFF` — cards, content panels
- Text Primary `#1A1A2E` — all headings and body text
- Text Secondary `#5A5A7A` — captions, supporting text

**Typography (1920×1080 stage)**
- Hero title: ExtraBold 800, 96–120px
- Slide title: SemiBold 600, 52–64px
- Body / bullets: Regular 400, 28–36px
- Label: SemiBold 600, 22px — ALL CAPS, 3px letter-spacing
- Stat number: ExtraBold 800, 120–160px

**Layout language**
- Title slides: left panel ~1200px (logo + label + title + subtitle); right panel ~720px (`#E8E6FF` bg, `clip-path:ellipse(100% 100% at 100% 50%)`, primary illustration)
- Content slides: MPG logo watermark top-right; white cards with 4px `#B8B4FF` left border; 60% text / 40% illustration
- Stats: white cards, oversized ExtraBold number in `#B8B4FF`, SemiBold label below
- Pricing / packages: two-col layout, 3 items auto-creates 3 columns; use `featured:true` on the best-value item

---

## Illustration Index

| File (no extension) | Use for… |
|---|---|
| `salary-transparency` | Pay gap, salary transparency, compensation, benchmarking — MPG core topic |
| `marketing` | Social media, content marketing, brand awareness, campaigns |
| `analytics` | Data, statistics, audience numbers, metrics, reach |
| `community-contribution` | Community, organic growth, grassroots movement |
| `human-resources` | HR, employer branding, recruitment, workplace |
| `event-management` | Events, seminars, conferences, meetups, activations |
| `financial-services` | Finance, investment, money, pricing |
| `education-ed-tech` | Education, learning, financial literacy, training |
| `consulting` | Consulting, advisory, strategy, proposals |
| `collaboration-communication` | Partnership, collaboration, B2B, joint ventures |
| `information-technology` | Tech, platforms, digital products, apps |
| `data-management` | Data management, databases, privacy, benchmarking tools |
| `sales` | Sales, conversion, revenue, business development |
| `product-management` | Product roadmap, features, platform updates |
| `project-management` | Project planning, timelines, deliverables |
| `accounting` | Accounting, bookkeeping, financial reporting |
| `legal` | Legal, compliance, contracts |
| `compliance` | Regulatory compliance, policies, standards |
| `customer-support` | Support, service, client management |
| `ai` | AI, automation, machine learning, tech innovation |
| `graphic-design` | Design, creative, visual content, branding |
| `professional-services` | Professional services, B2B overview |
| `ecommerce` | E-commerce, online retail, digital sales |
| `retail` | Retail, consumer, shopping |
| `healthcare` | Health, wellness, healthcare |
| `hospitality` | Hospitality, tourism, F&B |
| `food-and-beverages` | F&B industry |
| `manufacturing` | Manufacturing, production, supply chain |
| `construction` | Construction, real estate, infrastructure |
| `oil-and-gas` | Energy, oil and gas |
| `coffee-break` | Breaks, informal settings, culture slides |
| `mindless-1` | Abstract decorative — closing or warm slides |
| `mindless-2` | Abstract figure — "about us", human interest |
| `mindless-3` | Abstract shape — section break or divider |

---

## MPG Core Data

Use these facts whenever generating slides about MPG's reach, audience, or track record. Do not invent numbers — use these.

**Platform reach**
- Instagram: **285K** followers
- TikTok: **145K** followers
- Threads + Meta: **65K** combined
- YouTube: **1.2M** total views
- Mailing list: **12K** subscribers
- Total organic reach: **~52M views/year** (no paid ads — entirely organic)

**Audience profile**
- Age: primarily **25–34** (Millennials + early Gen Z)
- **60%+** work at MNCs or large corporates
- Geographic: Klang Valley, Johor, Penang
- High-intent: salary negotiators, career switchers, DEI advocates, HR leads

**Proof points / case studies**
- Open Company × CoinGecko: **451K** YouTube impressions, **139K** Instagram views — content-only, no live event
- Open Company has run **2+ years** with corporate partners
- MPG produced **salary benchmarks** used by thousands of Malaysian job seekers

---

## Deck Type Templates

Each template shows: slide number, suggested **layout**, label, content notes.

### 1. Partner / Sponsorship Proposal

1. **Cover** — `title` — heading: "[Partner] × MalaysianPAYGAP", subheading: deck tagline or date, illustration: `collaboration-communication`
2. **About MPG** — `content` — label: "WHO WE ARE", heading: "Closing Malaysia's *Pay Gap*", body: mission (1 sentence), bullets: 3 key facts (e.g. community size, what MPG does, why it matters), illustration: `salary-transparency`
3. **Our Audience** — `stats` — label: "OUR COMMUNITY", stats: age (25–34), MNC share (60%+), geography (3 cities), body: "High-purchase-intent professionals — salary negotiators, DEI champions, career switchers.", illustration: `community-contribution`
4. **Our Numbers** — `stats` — label: "OUR REACH", stats: 285K Instagram, 145K TikTok, 52M views/year, 12K mailing list, body: "Entirely organic. Zero paid ads.", illustration: `analytics`
5. **The Opportunity** — `content` — label: "THE BRIEF", heading: "What [Partner] Wants to *Achieve*", body: partner's goals (from document), bullets: 3–4 desired outcomes
6. **What We Propose** — `content` — label: "OUR PROPOSAL", heading: "A *Tailored* Partnership Plan", bullets: 4–5 deliverables (content types, event, timeline), illustration: `project-management`
7. **Why It Works** — `stats` — label: "PROOF OF CONCEPT", heading: "Open Company × *CoinGecko*", stats: 451K YouTube impressions, 139K Instagram views, body: "That was content alone — no live event. Open Company adds the in-person layer on top.", illustration: `marketing`
8. **Packages & Pricing** — `two-col` — label: "INVESTMENT", heading: "Choose Your *Level*", bullets: 3 packages, `**RM X,XXX**` at start of text, `featured:true` on best value, full deliverable list in body
9. **Next Steps** — `content` — label: "LET'S MOVE FORWARD", heading: "What *Happens* Next", bullets: 3–4 action steps (e.g. "Sign MOU by X", "Align on content brief", "Launch announcement"), illustration: `collaboration-communication`

### 2. Media Kit / Rate Card

1. **Cover** — `title` — heading: "MalaysianPAYGAP *Media Kit*", subheading: current year, illustration: `marketing`
2. **What Our Community Cares About** — `content` — label: "OUR AUDIENCE'S WORLD", heading: "They Want *Pay Equity* — And They're Talking About It", bullets: 4 topics (salary transparency, workplace fairness, equal pay, negotiation), illustration: `salary-transparency`
3. **Audience Snapshot** — `table` — label: "WHO'S LISTENING", heading: "The *Demographics*", rows: Age (25–34), Gender split, Location (KL / JB / Penang), Sector (60%+ MNCs), Intent (job seekers + HR leads)
4. **Our Reach** — `stats` — label: "OUR NUMBERS", heading: "**52M** Views. *Organically.*", stats: 285K Instagram, 145K TikTok, 1.2M YouTube views, 12K email, body: "No paid boosting. Every view is earned through content people actually want.", illustration: `analytics`
5. **Organic Credibility** — `content` — label: "WHY THIS MATTERS", heading: "Organic Reach Converts *Better*", body: "Paid campaigns interrupt. Organic communities trust. Our audience follows MPG because they want to — that trust extends to brands we partner with.", bullets: 3 engagement proof points, illustration: `community-contribution`
6. **How Brands Come In** — `two-col` — label: "PARTNERSHIP FORMATS", heading: "How We *Collaborate*", bullets: 3–4 options (Sponsored content, Event sponsorship, Content series, Social media takeover) with body describing each
7. **Packages & Rates** — `two-col` — label: "INVESTMENT", heading: "Pricing *Packages*", 3 packages with `**RM X,XXX**` and full deliverable lists
8. **Contact / Let's Talk** — `content` — label: "GET IN TOUCH", heading: "Ready to *Partner*?", body: contact info, illustration: `collaboration-communication`

### 3. Event Proposal / Brief

1. **Cover** — `title` — heading: event name, subheading: date + venue or "Open Company by MalaysianPAYGAP", illustration: `event-management`
2. **Event Concept** — `content` — label: "THE CONCEPT", heading: "What is *Open Company*?", body: format description, bullets: 3–4 core elements (venue, format, why unique), illustration: `event-management`
3. **Target Audience** — `content` — label: "WHO ATTENDS", heading: "Built for *Decision Makers*", bullets: 4–5 audience segments, body: why they come, illustration: `human-resources`
4. **Past Event Recaps** — `stats` — label: "IT ALREADY WORKS", stats: attendee count, engagement metrics, number of editions run, body: testimonial or key insight, illustration: `analytics`
5. **Event Format** — `table` — label: "THE DETAILS", heading: "What a *Typical* Edition Looks Like", rows: Venue, Audience (cap), Duration, Session format, Language, Frequency
6. **What We Propose** — `content` — label: "THIS EDITION", heading: "Our *Proposal* for [Partner]", bullets: programme outline (office tour → session → Q&A → networking), body: what partner gains, illustration: `project-management`
7. **Marketing & Promotion Plan** — `content` — label: "HOW WE DRIVE ATTENDANCE", heading: "From *Invitation* to Sold-Out", bullets: 4 channels (Instagram, email, Threads, direct outreach), illustration: `marketing`
8. **Co-branding & Sponsorship** — `two-col` — label: "VISIBILITY FOR PARTNERS", heading: "What *You* Get", bullets: 2–3 visibility tiers with descriptions
9. **Budget & Timeline** — `table` — label: "INVESTMENT & TIMELINE", heading: "Making It *Happen*", rows: dates, budget items, deliverables, owner
10. **Next Steps** — `content` — label: "LET'S DO THIS", heading: "Next *Steps*", bullets: 3–4 actions

### 4. Internal / Strategy Deck

1. **Cover** — `title` — heading: project/initiative name, subheading: team + date, illustration: `consulting`
2. **Context / Background** — `content` — label: "THE CONTEXT", heading: "Why This *Matters*", body: problem statement, bullets: 2–3 driving factors, illustration: `data-management`
3. **Goals** — `content` — label: "OUR GOALS", heading: "What We Want to *Achieve*", bullets: 3–5 goals with success metrics (use `→` style), illustration: `product-management`
4. **Approach** — `content` — label: "HOW WE'LL DO IT", heading: "The *Plan*", bullets: 3–5 steps or workstreams, illustration: `project-management`
5. **Timeline / Milestones** — `table` — label: "TIMELINE", heading: "*Milestone* Plan", rows: phase names with dates and owners
6. **Resources Needed** — `two-col` — label: "WHAT WE NEED", heading: "Resources & *Requirements*", bullets: budget, people, tools, dependencies
7. **Expected Outcomes** — `stats` — label: "EXPECTED IMPACT", heading: "What *Success* Looks Like", stats: 3 key metrics, body: summary of impact
8. **Next Steps / Owner** — `content` — label: "NEXT STEPS", heading: "Who Does *What*", bullets: 3–4 named actions with owners, illustration: `collaboration-communication`

---

## Layout Design Rules

**Title slide** — NON-NEGOTIABLE two-panel layout: left panel logo + label + heading; right panel `#E8E6FF` with illustration. Co-branded covers: `[Partner] × MalaysianPAYGAP` in heading with `×` in purple.

**Content slides** — 60% left for text, 40% right for illustration. Every content slide needs: label (uppercase) + heading (with *accent*) + either body card OR bullets (not both unless content demands it) + illustration when topic-matched.

**Stats slides** — 2–4 stat cards max. Always add `body` with the "so what" explanation. Always add `illustration`. Stat numbers must include units or K/M suffixes.

**Table slides** — use for any structured data: event formats, programme schedules, comparison specs, terms and conditions. 4–6 rows is optimal. Row label in purple, value in dark text.

**Two-col slides** — for packages, options, comparisons. 3 packages = 3 columns. Pricing cards must have `**RM X,XXX**` at start of text so the price renders large. Body = bulleted deliverable list (newline-separated).

**Checklist / goals slides** — use `content` layout with bullets. Format each item as the goal with supporting detail in `body`.

**Density**:
- Proposal/pitch decks: 1 main idea per slide, generous space, max 4 bullets, prefer illustrations
- Reading/send-ahead decks: use tables and body cards, 5–6 bullets OK, more detail per slide

---

## Intent Detection

If the user asks a question, makes a comment, or says anything that does NOT require changing slides → respond ONLY with:
`[{{"op":"CHAT","message":"your answer"}}]`

CHAT triggers: "which model?", "how many slides?", "what can you do?", "looks good", "what is MPG?", "explain this"
Deck op triggers: "create a deck", "add a slide", "make the title shorter", "generate from this PDF", any content edit

Never mix CHAT with deck operations.

---

## Valid Operations

{{"op":"CHAT","message":"..."}}                                          — conversational reply, no deck change
{{"op":"REPLACE_DECK","title":"...","slides":[...]}}                     — replace whole deck (fresh generation)
{{"op":"ADD_SLIDE","layout":"content","after":-1,"blocks":[...]}}
{{"op":"UPDATE_BLOCK","slideId":"s1","blockId":"b1","updates":{{"text":"..."}}}}
{{"op":"DELETE_SLIDE","slideId":"s1"}}
{{"op":"ADD_BLOCK","slideId":"s1","block":{{...}}}}

---

## Layouts

- **title** — two-panel: LEFT = logo + label + heading + subheading; RIGHT = illustration on `#E8E6FF`
- **content** — label + heading + body card and/or bullets + optional illustration
- **two-col** — label + heading + card grid. Item count drives columns: 2→2col, 3→3col, 4→4col. Items support `featured:true` for a Best Value highlight card.
- **stats** — label + heading + 2–4 stat cards (large number + label) + optional `body` block below + optional `illustration` block in right column
- **table** — label + heading + rows of key–value pairs (use `bullets` items where `text` = row label, `body` = row value); optional illustration in right column

---

## Block Types

- `label`       : `{{type, text}}` — uppercase section label
- `heading`     : `{{type, text}}` — main title; supports **bold** and *accent* markers
- `subheading`  : `{{type, text}}` — subtitle / tagline
- `body`        : `{{type, text}}` — paragraph card; supports **bold** markers
- `bullets`     : `{{type, items:[{{text, body?, featured?}}]}}` — bullet list or card grid
- `stats`       : `{{type, stats:[{{num, label}}]}}` — e.g. `{{"num":"**72%**","label":"No pay policy"}}`
- `illustration`: `{{type, asset:"<filename>"}}` — from index above (no file extension)
- `logo`        : `{{type, asset:"mpg_purple_logo"}}` — top-left on title; watermark on content

---

## Text Emphasis

- `**word**` → bold accent-purple. Use for: key stats, prices, brand names, impact numbers.
- `*word*`   → accent-purple non-bold. Use for: one key concept word per heading.

Examples:
  heading:       `"The *Gender* Pay Gap"` / `"Why Companies *Partner* With Us"`
  body:          `"Recruiters charge **15–25%** of salary. We charge nothing."`
  stat:          `{{"num":"**72%**","label":"Companies with no pay policy"}}`
  pricing card:  `{{"text":"**RM 30,000** Full Package","featured":true,"body":"Everything in Social + Video Bundle\nLong-form YouTube documentary\nShort reel (IG + TikTok)\nEvent coordination by MPG\n+ RM 80 / participant"}}`
  table row:     `{{"text":"Audience","body":"~40 participants — capped for quality"}}`

---

## Content Rules

1. Every content/two-col/stats/table slide MUST have a `label` block (e.g. "THE PROBLEM", "OUR REACH").
2. Headings: always use `*accent*` on the key concept word.
3. **Pricing cards** (`two-col`): put the price as `**RM X,XXX**` at the START of the `text` field, followed by the package name. Example: `"text": "**RM 30,000** Full Package"`. The body field must list ALL deliverables — minimum 3 lines separated by `\n`. Mark the best-value package with `"featured": true`.
4. **Stats slides**: always add a `body` block with a 1–2 sentence explanation of what the numbers mean. Add an `illustration` block matching the topic.
5. **Table layout**: use when presenting structured details (event format, programme schedule, terms, specs). Each `bullets` item: `text` = row label, `body` = row value. Include 4–6 rows minimum.
6. Stats: use real numbers with units. Format numbers as `**X%**` or `**451K**` in the `num` field.
7. Title slide: always include an `illustration` block and a `logo` block.
8. Target 10–14 slides for a full deck. Rich content preferred — do not skip details. Every slide should be fully populated.
9. Max 100 words per slide. Aim for at least 50 words of visible content on every non-title slide.

---

## Available Brand Assets

{brand_assets}

---

Respond with ONLY a valid JSON array. No explanation, no markdown fences, no prose.
"""

_SYSTEM_PROMPT = _load_system_prompt()


# ── Request models ─────────────────────────────────────────────────────────────
class GenerateRequest(BaseModel):
    message: str
    deck: dict
    session_id: Optional[str] = None
    model_id: Optional[str] = None      # override model name
    endpoint_id: Optional[str] = None   # override endpoint (ensures URL+headers match)


class ExportPptxRequest(BaseModel):
    deck: dict
    brand: Optional[str] = "mpg"


class PaletteUpdateRequest(BaseModel):
    palette: list


# ── Route factory ──────────────────────────────────────────────────────────────
def setup_deck_builder_routes() -> APIRouter:
    router = APIRouter()

    # ── POST /api/deck/extract-pdf ─────────────────────────────────────────────
    @router.post("/api/deck/extract-pdf")
    async def extract_pdf(request: Request, file: UploadFile = File(...)):
        owner = require_user(request)
        if not file.filename or not file.filename.lower().endswith(".pdf"):
            raise HTTPException(400, "Only PDF files are accepted")

        content = await file.read(10 * 1024 * 1024)  # 10 MB cap
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            from src.personal_docs import extract_pdf_text
            text = extract_pdf_text(tmp_path)
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass

        if not text.strip():
            raise HTTPException(422, "Could not extract text from this PDF")

        return {"text": text[:50000], "chars": len(text)}

    # ── POST /api/deck/generate ────────────────────────────────────────────────
    @router.post("/api/deck/generate")
    async def generate_deck(request: Request, body: GenerateRequest):
        owner = require_user(request)
        from src.endpoint_resolver import resolve_endpoint
        from src.llm_core import llm_call_async
        from src.text_helpers import strip_think

        url, model, headers = resolve_endpoint("utility", owner=owner or None)
        if not url:
            url, model, headers = resolve_endpoint("default", owner=owner or None)
        if not url or not model:
            raise HTTPException(503, "No LLM endpoint configured")
        # If the frontend picked a specific endpoint+model, resolve it properly
        # (avoids sending a claude-* model name to a deepseek endpoint URL)
        if body.endpoint_id:
            from src.endpoint_resolver import resolve_endpoint_by_id
            result = resolve_endpoint_by_id(body.endpoint_id, body.model_id, owner=owner or None)
            if result:
                url, model, headers = result
        elif body.model_id:
            model = body.model_id

        deck_summary = json.dumps(body.deck, ensure_ascii=False)[:8000]
        user_msg = f"{body.message}\n\nCurrent deck state:\n{deck_summary}"

        try:
            raw = await llm_call_async(
                url=url,
                model=model,
                messages=[
                    {"role": "system", "content": _SYSTEM_PROMPT},
                    {"role": "user",   "content": user_msg},
                ],
                headers=headers,
                temperature=0.3,
                max_tokens=int(get_setting("deck_max_tokens", 6000)),
                timeout=120,
            )
        except Exception as e:
            logger.error("deck generate LLM call failed: %s", e)
            raise HTTPException(502, f"LLM call failed: {e}")

        cleaned = strip_think(raw or "", prose=False, prompt_echo=True)
        cleaned = re.sub(r"^```(?:json)?\s*|\s*```$", "", cleaned, flags=re.MULTILINE).strip()
        match = re.search(r"\[[\s\S]*\]", cleaned)
        if not match:
            raise HTTPException(422, detail={"error": "LLM did not return a JSON array", "raw": cleaned[:400]})

        try:
            ops = json.loads(match.group(0))
        except json.JSONDecodeError as e:
            raise HTTPException(422, detail={"error": f"JSON parse error: {e}", "raw": cleaned[:400]})

        if not isinstance(ops, list):
            raise HTTPException(422, "Expected a JSON array of operations")

        return {"operations": ops}

    # ── POST /api/deck/export-pptx ─────────────────────────────────────────────
    @router.post("/api/deck/export-pptx")
    async def export_pptx(request: Request, body: ExportPptxRequest):
        owner = require_user(request)
        try:
            from pptx import Presentation
            from pptx.util import Emu, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise HTTPException(503, "python-pptx is not installed. Run: pip install python-pptx")

        deck = body.deck
        brand_id = body.brand or "mpg"

        # Load brand tokens
        brand_tokens = _load_brand_tokens(brand_id)

        prs = Presentation()
        prs.slide_width  = Emu(12_192_000)   # 13.333 inches
        prs.slide_height = Emu(6_858_000)    # 7.5 inches

        slides = deck.get("slides") or []
        if not slides:
            raise HTTPException(422, "Deck has no slides")

        blank_layout = prs.slide_layouts[6]  # blank layout

        for slide_data in slides:
            pslide = prs.slides.add_slide(blank_layout)

            # Background fill
            bg_color_hex = (
                (slide_data.get("background") or {}).get("color")
                or brand_tokens.get("bg", "#F2F2F2")
            )
            fill = pslide.background.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(bg_color_hex)

            _render_slide_pptx(pslide, slide_data, brand_tokens)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        title = re.sub(r"[^\w\s-]", "", deck.get("title") or "presentation")
        title = re.sub(r"\s+", "-", title.strip()).lower() or "deck"
        filename = f"{title}.pptx"

        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    # ── GET /api/brand/kit ─────────────────────────────────────────────────────
    @router.get("/api/brand/kit")
    async def get_brand_kit(request: Request, brand: str = "mpg"):
        owner = require_user(request)
        brand_dir = _BRAND_DIR / brand
        if not brand_dir.exists():
            raise HTTPException(404, f"Brand '{brand}' not found")

        brand_json = _load_brand_tokens(brand)

        # List asset categories
        assets = {}
        for category in ("logos", "illustrations"):
            cat_dir = brand_dir / category
            if cat_dir.exists():
                assets[category] = [
                    {"name": f.stem, "filename": f.name, "url": f"/api/brand/asset/{brand}/{category}/{f.name}"}
                    for f in sorted(cat_dir.iterdir())
                    if f.is_file() and f.suffix.lower() in {".png", ".svg", ".jpg", ".webp"}
                ]
            else:
                assets[category] = []

        return {"brand": brand_json, "assets": assets}

    # ── GET /api/brand/font-css/{brand} ───────────────────────────────────────
    @router.get("/api/brand/font-css/{brand}")
    async def get_font_css(request: Request, brand: str = "mpg"):
        """Return base64-embedded @font-face CSS for the brand's Poppins fonts."""
        owner = require_user(request)
        import base64
        fonts_dir = _BRAND_DIR / brand / "fonts"
        font_map = [
            ("poppins-regular.ttf",         "400", "normal"),
            ("poppins-italic.ttf",          "400", "italic"),
            ("poppins-medium.ttf",          "500", "normal"),
            ("poppins-semibold.ttf",        "600", "normal"),
            ("poppins-bold.ttf",            "700", "normal"),
            ("poppins-extrabold.ttf",       "800", "normal"),
            ("poppins-extrabolditalic.ttf", "800", "italic"),
        ]
        parts = []
        for filename, weight, style in font_map:
            fp = fonts_dir / filename
            if fp.exists():
                b64 = base64.b64encode(fp.read_bytes()).decode()
                parts.append(
                    f"@font-face{{font-family:'Poppins';"
                    f"src:url('data:font/truetype;base64,{b64}') format('truetype');"
                    f"font-weight:{weight};font-style:{style};font-display:swap;}}"
                )
        css = "\n".join(parts) if parts else (
            "@import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700;800&display=swap');"
        )
        return Response(content=css, media_type="text/css")

    # ── GET /api/brand/asset/{brand}/{category}/{filename} ─────────────────────
    @router.get("/api/brand/asset/{brand}/{category}/{filename}")
    async def serve_brand_asset(request: Request, brand: str, category: str, filename: str):
        owner = require_user(request)
        # Safety: no path traversal
        if ".." in brand or ".." in category or ".." in filename:
            raise HTTPException(400, "Invalid path")
        asset_path = _BRAND_DIR / brand / category / filename
        if not asset_path.exists():
            raise HTTPException(404, "Asset not found")

        suffix = asset_path.suffix.lower()
        media_types = {".png": "image/png", ".svg": "image/svg+xml", ".jpg": "image/jpeg", ".webp": "image/webp"}
        media_type = media_types.get(suffix, "application/octet-stream")

        return StreamingResponse(open(asset_path, "rb"), media_type=media_type)

    # ── POST /api/brand/upload ─────────────────────────────────────────────────
    @router.post("/api/brand/upload")
    async def upload_brand_asset(
        request: Request,
        file: UploadFile = File(...),
        brand: str = Form("mpg"),
        category: str = Form("illustrations"),
    ):
        owner = require_user(request)
        if category not in ("logos", "illustrations"):
            raise HTTPException(400, "category must be 'logos' or 'illustrations'")
        if file.content_type not in _ALLOWED_IMAGE_TYPES:
            raise HTTPException(415, "Only PNG, JPEG, SVG, and WebP are accepted")

        content = await file.read(5 * 1024 * 1024)  # 5 MB cap
        dest_dir = _BRAND_DIR / brand / category
        dest_dir.mkdir(parents=True, exist_ok=True)

        safe_name = re.sub(r"[^\w.\-]", "_", Path(file.filename or "asset").name)
        dest_path = dest_dir / safe_name
        dest_path.write_bytes(content)

        return {
            "ok": True,
            "name": dest_path.stem,
            "filename": dest_path.name,
            "url": f"/api/brand/asset/{brand}/{category}/{dest_path.name}",
        }

    # ── DELETE /api/brand/asset/{brand}/{category}/{name} ─────────────────────
    @router.delete("/api/brand/asset/{brand}/{category}/{name}")
    async def delete_brand_asset(request: Request, brand: str, category: str, name: str):
        owner = require_user(request)
        if ".." in brand or ".." in category or ".." in name:
            raise HTTPException(400, "Invalid path")
        # name may or may not include extension — try both
        base = _BRAND_DIR / brand / category
        deleted = False
        for f in base.iterdir() if base.exists() else []:
            if f.stem == name or f.name == name:
                f.unlink()
                deleted = True
                break
        if not deleted:
            raise HTTPException(404, "Asset not found")
        return {"ok": True}

    # ── PUT /api/brand/palette ─────────────────────────────────────────────────
    @router.put("/api/brand/palette")
    async def update_palette(request: Request, body: PaletteUpdateRequest, brand: str = "mpg"):
        owner = require_user(request)
        brand_file = _BRAND_DIR / brand / "brand.json"
        if not brand_file.exists():
            raise HTTPException(404, f"Brand '{brand}' not found")
        data = json.loads(brand_file.read_text(encoding="utf-8"))
        data["palette"] = body.palette
        brand_file.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        return {"ok": True, "palette": data["palette"]}

    return router


# ── Helpers ────────────────────────────────────────────────────────────────────

def _load_brand_tokens(brand_id: str) -> dict:
    brand_file = _BRAND_DIR / brand_id / "brand.json"
    if brand_file.exists():
        try:
            return json.loads(brand_file.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {}


def _hex_to_rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return RGBColor(0x1A, 0x1A, 0x2E)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def _pct_to_emu(pct: float, total_emu: int) -> int:
    return int(max(0.0, min(1.0, pct)) * total_emu)


_SLIDE_W = 12_192_000   # 13.333 in @ 914400 EMU/in
_SLIDE_H =  6_858_000   #  7.5   in


def _emu(x_pct: float, y_pct: float, w_pct: float, h_pct: float):
    """(fraction × 4) → (emu × 4) clamped to slide bounds."""
    def _c(v, tot): return int(max(0.0, min(1.0, v)) * tot)
    return _c(x_pct, _SLIDE_W), _c(y_pct, _SLIDE_H), _c(w_pct, _SLIDE_W), _c(h_pct, _SLIDE_H)


def _strip_md(text: str) -> str:
    """Remove **bold** and *accent* markdown markers → plain text."""
    s = re.sub(r'\*\*([^*\n]+?)\*\*', r'\1', str(text or ''))
    return re.sub(r'\*([^*\n]+?)\*', r'\1', s)


def _md_runs(para, text: str, pt: float, color_hex: str, bold: bool, accent_hex: str) -> None:
    """Add styled runs: **x** → bold accent colour, *x* → accent colour."""
    from pptx.util import Pt
    text = str(text or '')
    pat = re.compile(r'\*\*([^*\n]+?)\*\*|\*([^*\n]+?)\*')
    last = 0
    for m in pat.finditer(text):
        if m.start() > last:
            r = para.add_run(); r.text = text[last:m.start()]
            r.font.size = Pt(pt); r.font.bold = bold
            r.font.color.rgb = _hex_to_rgb(color_hex); r.font.name = 'Poppins'
        r = para.add_run()
        r.text = m.group(1) if m.group(1) is not None else m.group(2)
        r.font.bold = m.group(1) is not None
        r.font.size = Pt(pt)
        r.font.color.rgb = _hex_to_rgb(accent_hex); r.font.name = 'Poppins'
        last = m.end()
    if last < len(text):
        r = para.add_run(); r.text = text[last:]
        r.font.size = Pt(pt); r.font.bold = bold
        r.font.color.rgb = _hex_to_rgb(color_hex); r.font.name = 'Poppins'


def _txb(pslide, x: int, y: int, w: int, h: int, wrap: bool = True):
    """Add a textbox and return its TextFrame."""
    from pptx.util import Emu
    box = pslide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    box.text_frame.word_wrap = wrap
    return box.text_frame


def _rect(pslide, x: int, y: int, w: int, h: int, fill_hex: str, border_hex: str = None):
    """Add a filled rectangle shape."""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    s = pslide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    if border_hex:
        s.line.color.rgb = _hex_to_rgb(border_hex)
        s.line.width = Emu(int(0.003 * _SLIDE_W))
    else:
        s.line.fill.background()
    return s


def _place_png(pslide, data_uri: str, x: int, y: int, w: int, h: int) -> bool:
    """Embed a data:image/... URI as a picture. Returns True on success."""
    import base64
    if not (data_uri or '').startswith('data:image/'):
        return False
    try:
        header, b64 = data_uri.split(',', 1)
        img_bytes = base64.b64decode(b64)
        sfx = '.png' if 'png' in header else '.jpg'
        with tempfile.NamedTemporaryFile(suffix=sfx, delete=False) as f:
            f.write(img_bytes); tmp = f.name
        try:
            from pptx.util import Emu
            pslide.shapes.add_picture(tmp, Emu(x), Emu(y), Emu(w), Emu(h))
        finally:
            try: os.unlink(tmp)
            except Exception: pass
        return True
    except Exception as e:
        logger.warning('picture embed: %s', e)
        return False


def _get_blk(blocks: list, btype: str):
    return next((b for b in blocks if (b.get('type') or '').lower() == btype), None)


def _render_slide_pptx(pslide, slide_data: dict, brand: dict) -> None:
    """Layout-aware slide renderer for PPTX export."""
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN

    layout  = (slide_data.get('layout') or 'content').lower()
    blocks  = slide_data.get('blocks') or []
    sd      = brand.get('slide_defaults') or {}
    accent  = sd.get('accent', '#B8B4FF')
    card    = sd.get('card',   '#E8E6FF')
    dark    = sd.get('text',   '#1A1A2E')
    muted   = '#5A5A7A'

    label_b = _get_blk(blocks, 'label')
    head_b  = _get_blk(blocks, 'heading')
    sub_b   = _get_blk(blocks, 'subheading')
    body_b  = _get_blk(blocks, 'body')
    bul_b   = _get_blk(blocks, 'bullets')
    stat_b  = _get_blk(blocks, 'stats')
    illus_b = _get_blk(blocks, 'illustration')
    logo_b  = _get_blk(blocks, 'logo')

    # ── inner helpers ─────────────────────────────────────────────────────────
    def lbl(x, y, w, h, text):
        tf = _txb(pslide, *_emu(x, y, w, h), wrap=False)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = _strip_md(text).upper()
        r.font.size = Pt(9); r.font.bold = True
        r.font.color.rgb = _hex_to_rgb(accent); r.font.name = 'Poppins'

    def hd(x, y, w, h, text, pt=30):
        tf = _txb(pslide, *_emu(x, y, w, h))
        _md_runs(tf.paragraphs[0], text, pt, dark, True, accent)

    def watermark():
        ld = _BRAND_DIR / 'mpg' / 'logos'
        for cand in ('mpg_purple_logo.png', 'mpg-purple-logo.png'):
            lp = ld / cand
            if lp.exists():
                lx, ly, lw, lh = _emu(0.84, 0.03, 0.12, 0.10)
                try: pslide.shapes.add_picture(str(lp), Emu(lx), Emu(ly), Emu(lw), Emu(lh))
                except Exception: pass
                return

    # ── TITLE ─────────────────────────────────────────────────────────────────
    if layout == 'title':
        _rect(pslide, *_emu(0.58, 0, 0.42, 1.0), card)
        ld = _BRAND_DIR / 'mpg' / 'logos'
        la = (logo_b or {}).get('asset') or 'mpg_purple_logo'
        for cand in (f'{la}.png', 'mpg_purple_logo.png', 'mpg-purple-logo.png'):
            lp = ld / cand
            if lp.exists():
                try:
                    llx, lly, llw, llh = _emu(0.04, 0.04, 0.18, 0.14)
                    pslide.shapes.add_picture(str(lp), Emu(llx), Emu(lly), Emu(llw), Emu(llh))
                except Exception: pass
                break
        if label_b: lbl(0.05, 0.28, 0.50, 0.07, label_b.get('text') or '')
        if head_b:  hd(0.05, 0.36, 0.50, 0.32, head_b.get('text') or '', pt=38)
        if sub_b:
            tf = _txb(pslide, *_emu(0.05, 0.72, 0.50, 0.14))
            _md_runs(tf.paragraphs[0], sub_b.get('text') or '', 15, muted, False, accent)
        if illus_b:
            _place_png(pslide, illus_b.get('png') or illus_b.get('src') or '',
                       *_emu(0.61, 0.08, 0.35, 0.84))
        return

    watermark()

    # ── TWO-COL ───────────────────────────────────────────────────────────────
    if layout in ('two-col', 'two_col', 'bullets-card', 'bullets_card'):
        if label_b: lbl(0.05, 0.05, 0.90, 0.08, label_b.get('text') or '')
        if head_b:  hd(0.05, 0.13, 0.90, 0.16, head_b.get('text') or '', pt=30)

        items = (bul_b or {}).get('items') or []
        cols  = min(max(len(items), 1), 4)
        GAP   = 0.025
        X0    = 0.04
        TW    = 0.92
        cw    = (TW - GAP * (cols - 1)) / cols
        CY, CH = 0.31, 0.62
        fs    = 0.85 if cols >= 3 else 1.0

        for ci, item in enumerate(items):
            raw  = (item.get('text') or str(item)) if isinstance(item, dict) else str(item)
            bd2  = (item.get('body') or '') if isinstance(item, dict) else ''
            feat = isinstance(item, dict) and bool(item.get('featured'))
            cx   = X0 + ci * (cw + GAP)

            lx, ly, lw, lh = _emu(cx, CY, cw, CH)
            _rect(pslide, lx, ly, lw, lh, card if feat else '#FFFFFF',
                  accent if feat else None)
            if not feat:
                _rect(pslide, lx, ly, int(0.005 * _SLIDE_W), lh, accent)

            if feat:
                bx, by, bw, bh = _emu(cx + cw * 0.15, CY - 0.028, cw * 0.70, 0.048)
                badge = _rect(pslide, bx, by, bw, bh, accent)
                btf = badge.text_frame; btf.word_wrap = False
                bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
                br = bp.add_run(); br.text = 'BEST VALUE'
                br.font.size = Pt(8); br.font.bold = True
                br.font.color.rgb = _hex_to_rgb('#FFFFFF'); br.font.name = 'Poppins'

            pm  = re.search(r'\*\*(RM[\s\d,.KkMmBb]+)\*\*', raw, re.IGNORECASE)
            price = pm.group(1).strip() if pm else None
            name  = re.sub(r'\*\*RM[^*]+\*\*', '', raw, flags=re.IGNORECASE).strip().lstrip('-–— ')

            ix  = cx + 0.015
            iw  = cw - 0.025
            cur = CY + 0.04

            if price:
                tf = _txb(pslide, *_emu(ix, cur, iw, CH * 0.28), wrap=False)
                r  = tf.paragraphs[0].add_run(); r.text = price
                r.font.size = Pt(int(32 * fs)); r.font.bold = True
                r.font.color.rgb = _hex_to_rgb(accent); r.font.name = 'Poppins'
                cur += CH * 0.28

            if name:
                hd(ix, cur, iw, CH * 0.11, name, pt=int(13 * fs))
                cur += CH * 0.11
                if price:
                    tf = _txb(pslide, *_emu(ix, cur, iw, CH * 0.07), wrap=False)
                    r  = tf.paragraphs[0].add_run(); r.text = 'per edition'
                    r.font.size = Pt(9)
                    r.font.color.rgb = _hex_to_rgb(muted); r.font.name = 'Poppins'
                    cur += CH * 0.07

            if bd2:
                lines = [l.strip() for l in re.split(r'\n|→|•', bd2) if l.strip()]
                rem = (CY + CH) - cur - 0.02
                if lines and rem > 0.05:
                    tf = _txb(pslide, *_emu(ix, cur, iw, rem))
                    for li, line in enumerate(lines):
                        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                        _md_runs(p, '→  ' + _strip_md(line), int(10 * fs), muted, False, accent)
        return

    # ── TABLE ─────────────────────────────────────────────────────────────────
    if layout == 'table':
        has_i = illus_b is not None
        TW    = 0.52 if has_i else 0.88
        if label_b: lbl(0.05, 0.05, TW, 0.08, label_b.get('text') or '')
        if head_b:  hd(0.05, 0.13, TW, 0.16, head_b.get('text') or '', pt=28)

        items = (bul_b or {}).get('items') or []
        n     = len(items)
        if n:
            RY0 = 0.31
            RH  = min(0.10, 0.62 / n)
            RG  = 0.006
            KW  = 0.16
            VX  = 0.05 + KW + 0.02
            VW  = TW - KW - 0.04
            for ri, item in enumerate(items):
                key = (item.get('text') or '') if isinstance(item, dict) else str(item)
                val = (item.get('body') or '') if isinstance(item, dict) else ''
                ry  = RY0 + ri * (RH + RG)
                _rect(pslide, *_emu(0.04, ry, TW - 0.01, RH),
                      '#FFFFFF' if ri % 2 == 0 else '#F4F4FA')
                tf = _txb(pslide, *_emu(0.05, ry + 0.01, KW, RH - 0.01), wrap=False)
                r  = tf.paragraphs[0].add_run(); r.text = _strip_md(key)
                r.font.size = Pt(11); r.font.bold = True
                r.font.color.rgb = _hex_to_rgb(accent); r.font.name = 'Poppins'
                tf = _txb(pslide, *_emu(VX, ry + 0.01, VW, RH - 0.01))
                _md_runs(tf.paragraphs[0], val, 11, dark, False, accent)

        if has_i:
            _place_png(pslide, illus_b.get('png') or illus_b.get('src') or '',
                       *_emu(0.59, 0.10, 0.37, 0.82))
        return

    # ── STATS ─────────────────────────────────────────────────────────────────
    if layout == 'stats':
        has_i = illus_b is not None
        TW    = 0.52 if has_i else 0.90
        if label_b: lbl(0.05, 0.05, TW, 0.08, label_b.get('text') or '')
        if head_b:  hd(0.05, 0.13, TW, 0.16, head_b.get('text') or '', pt=28)

        stats = (stat_b or {}).get('stats') or []
        n     = max(len(stats), 1)
        GAP   = 0.02
        X0    = 0.04
        SY    = 0.31
        SH    = 0.38 if body_b else 0.50
        sw    = (TW - 0.01 - GAP * (n - 1)) / n

        for si, stat in enumerate(stats):
            sx = X0 + si * (sw + GAP)
            _rect(pslide, *_emu(sx, SY, sw, SH), '#FFFFFF')
            num_pt = 40 if n >= 4 else (48 if n == 3 else 56)
            tf = _txb(pslide, *_emu(sx, SY + 0.03, sw, SH * 0.52), wrap=False)
            p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            _md_runs(p, _strip_md(str(stat.get('num') or '')), num_pt, accent, True, accent)
            tf = _txb(pslide, *_emu(sx, SY + SH * 0.58, sw, SH * 0.36))
            p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            _md_runs(p, _strip_md(str(stat.get('label') or '')), 12, dark, False, accent)

        if body_b:
            BY = SY + SH + 0.03
            bx, by, bw, bh = _emu(0.04, BY, TW - 0.01, max(1.0 - BY - 0.04, 0.08))
            _rect(pslide, bx, by, bw, bh, card)
            tf = _txb(pslide, bx, by, bw, bh)
            _md_runs(tf.paragraphs[0], body_b.get('text') or '', 13, dark, False, accent)

        if has_i:
            _place_png(pslide, illus_b.get('png') or illus_b.get('src') or '',
                       *_emu(0.59, 0.10, 0.37, 0.82))
        return

    # ── CONTENT (default) ─────────────────────────────────────────────────────
    has_i = illus_b is not None
    TW    = 0.52 if has_i else 0.90

    if label_b: lbl(0.05, 0.05, TW, 0.08, label_b.get('text') or '')
    if head_b:  hd(0.05, 0.13, TW, 0.18, head_b.get('text') or '', pt=30)

    cur = 0.33
    if sub_b and not body_b:
        tf = _txb(pslide, *_emu(0.05, cur, TW, 0.12))
        _md_runs(tf.paragraphs[0], sub_b.get('text') or '', 15, muted, False, accent)
        cur += 0.14

    if body_b:
        bx, by, bw, bh = _emu(0.05, cur, TW, 0.20)
        _rect(pslide, bx, by, bw, bh, '#FFFFFF')
        _rect(pslide, bx, by, int(0.004 * _SLIDE_W), bh, accent)
        tf = _txb(pslide, *_emu(0.065, cur + 0.015, TW - 0.025, 0.17))
        _md_runs(tf.paragraphs[0], body_b.get('text') or '', 14, dark, False, accent)
        cur += 0.22

    if bul_b:
        items = bul_b.get('items') or []
        rem   = max(1.0 - cur - 0.04, 0.10)
        tf    = _txb(pslide, *_emu(0.05, cur, TW, rem))
        for bi, item in enumerate(items):
            it   = (item.get('text') or str(item)) if isinstance(item, dict) else str(item)
            ib   = (item.get('body') or '') if isinstance(item, dict) else ''
            para = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            _md_runs(para, '→  ' + _strip_md(it), 13, dark, False, accent)
            if ib:
                sp = tf.add_paragraph(); sp.level = 1
                _md_runs(sp, '    ' + _strip_md(ib), 11, muted, False, accent)

    if has_i:
        _place_png(pslide, illus_b.get('png') or illus_b.get('src') or '',
                   *_emu(0.59, 0.08, 0.37, 0.84))
